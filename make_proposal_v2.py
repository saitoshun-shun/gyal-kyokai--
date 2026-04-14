"""
事業提案書PPT v2 - 最終版
ギャル協会×サイバーバズ×ENTIAL
米国TikTok×インバウンド体験×1億→10億
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PINK   = RGBColor(0xFF,0x69,0xB4)
LPINK  = RGBColor(0xFF,0xD6,0xE8)
ACCENT = RGBColor(0xFF,0x14,0x93)
DARK   = RGBColor(0x2D,0x2D,0x2D)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GRAY   = RGBColor(0x88,0x88,0x88)
BLACK  = RGBColor(0x00,0x00,0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, bg_color=None, border_color=None, bw=0):
    from pptx.util import Pt as P
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.background() if not bg_color else (sh.fill.solid(), setattr(sh.fill.fore_color,'rgb',bg_color))
    if bg_color: sh.fill.solid(); sh.fill.fore_color.rgb = bg_color
    else: sh.fill.background()
    if border_color and bw: sh.line.color.rgb = border_color; sh.line.width = P(bw)
    else: sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=16, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.italic = italic
    return tb

def hbar(slide, title, sub=None):
    box(slide,0,0,13.33,1.2,bg_color=PINK)
    txt(slide,title,0.4,0.12,11,0.75,size=30,bold=True,color=WHITE)
    if sub: txt(slide,sub,0.4,0.82,11,0.35,size=13,color=WHITE)

def num(slide,n):
    txt(slide,str(n),12.7,7.1,0.5,0.3,size=11,color=GRAY,align=PP_ALIGN.RIGHT)

# ═══════ SLIDE 1 表紙 ═══════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,LPINK)
box(s,0,2.0,13.33,3.4,bg_color=PINK)
txt(s,"ギャル × インバウンド",0.5,2.1,12.3,1.1,
    size=52,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"コンテンツビジネス事業提案書",0.5,3.15,12.3,0.8,
    size=26,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"ギャル協会  ×  サイバーバズ  ×  ENTIAL",
    0.5,5.6,12.3,0.6,size=20,bold=True,color=ACCENT,align=PP_ALIGN.CENTER)
txt(s,"Year1：1億円　／　Year3：10億円",
    0.5,6.25,12.3,0.45,size=16,color=DARK,align=PP_ALIGN.CENTER)
txt(s,"2026年4月",0.5,6.85,12.3,0.35,size=13,color=GRAY,align=PP_ALIGN.CENTER)

# ═══════ SLIDE 2 事業概要 ════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"事業概要・ビジョン"); num(s,2)

box(s,0.5,1.4,12.3,1.0,bg_color=LPINK,border_color=PINK,bw=2)
txt(s,"「ギャル文化」を米国TikTok×インバウンド体験で世界に届け、1年1億・3年10億を目指す",
    0.7,1.52,12.0,0.75,size=19,bold=True,color=ACCENT,align=PP_ALIGN.CENTER)

points = [
    ("メディア", "サイバーバズが米国向けTikTokアカウント（WESELL）を運営。ショートドラマ×グッズ販売で収益化"),
    ("体験",     "ppgalclub参考の渋谷ギャル体験をインバウンド商品化。外国人に「本物のギャル体験」を提供"),
    ("営業",     "ENTIALがバズ社内PMとして稼働。人件費ゼロで企業・自治体へのB2B営業を推進"),
    ("IP",       "うさたにパイセン／ギャル協会がコンテンツを創出。バズがメディア展開・WESELL販売で収益化"),
]
for i,(icon,desc) in enumerate(points):
    l = 0.4+i*3.2
    box(s,l,2.6,3.0,4.5,bg_color=LPINK if i%2==0 else WHITE,border_color=PINK,bw=1.5)
    box(s,l,2.6,3.0,0.65,bg_color=PINK)
    txt(s,icon,l+0.05,2.63,2.9,0.58,size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,desc,l+0.15,3.4,2.7,3.5,size=14,color=DARK,wrap=True)

# ═══════ SLIDE 3 3社の役割 ══════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"3社の役割　〜それぞれが何をするか〜"); num(s,3)

roles = [
    ("ギャル協会\nうさたにパイセン",
     "IPオーナー・\nコンテンツクリエイター",
     ["ショートドラマ出演・ネタ提供","ギャル文化の世界観発信","渋谷体験コンテンツの監修","海外向けコンテンツ創出"]),
    ("サイバーバズ",
     "メディアオーナー・\n資金提供・実行",
     ["TikTok/IGアカウント運営","WESELL（EC）で商品販売","制作費・活動費を全額負担","広告枠の販売・拡散"]),
    ("ENTIAL",
     "バズ社内PM・\n営業統括",
     ["バズ社内にPMとして常駐","企業・自治体への営業","営業代理店の開拓・管理","バズの人件費ゼロで稼働"]),
]
for i,(name,role,bullets) in enumerate(roles):
    l = 0.4+i*4.3
    box(s,l,1.35,4.0,5.75,bg_color=LPINK if i==1 else WHITE,border_color=PINK,bw=2)
    box(s,l,1.35,4.0,1.0,bg_color=PINK)
    txt(s,name,l+0.1,1.38,3.8,0.62,size=16,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,role,l+0.1,2.42,3.8,0.65,size=13,bold=True,color=ACCENT,align=PP_ALIGN.CENTER)
    for j,b in enumerate(bullets):
        txt(s,"▸  "+b,l+0.2,3.2+j*0.85,3.6,0.75,size=13,color=DARK)
    if i==2:
        box(s,l,6.6,4.0,0.45,bg_color=ACCENT)
        txt(s,"★ バズの人件費ゼロで稼働",l+0.1,6.63,3.8,0.38,
            size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

txt(s,"→",4.25,3.8,0.5,0.5,size=28,bold=True,color=PINK,align=PP_ALIGN.CENTER)
txt(s,"→",8.55,3.8,0.5,0.5,size=28,bold=True,color=PINK,align=PP_ALIGN.CENTER)

# ═══════ SLIDE 4 プラットフォーム戦略 ══════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"プラットフォーム戦略　〜TikTok優先×IGサブ〜"); num(s,4)

# TikTok（メイン）
box(s,0.4,1.4,6.0,5.6,bg_color=LPINK,border_color=PINK,bw=2)
box(s,0.4,1.4,6.0,0.75,bg_color=PINK)
txt(s,"TikTok（メイン）× WESELL",0.5,1.45,5.8,0.62,
    size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
tiktok_pts = [
    "ターゲット：米国（1.7億ユーザー）",
    "WESELLの切り抜き機能でコンテンツ→販売がワンステップ",
    "バズがWESELLノウハウを既に保有",
    "ショートドラマ → コスメ・グッズ直販",
    "英語コンテンツで全世界にリーチ",
    "フォロワー目標：Year1→1万、Year3→50万",
]
for i,pt in enumerate(tiktok_pts):
    txt(s,"▸  "+pt,0.6,2.35+i*0.78,5.6,0.68,size=13,color=DARK)

# Instagram（サブ）
box(s,6.9,1.4,6.0,5.6,bg_color=WHITE,border_color=PINK,bw=2)
box(s,6.9,1.4,6.0,0.75,bg_color=ACCENT)
txt(s,"Instagram（サブ）× AI翻訳",7.0,1.45,5.8,0.62,
    size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
ig_pts = [
    "TikTokで当たったコンテンツを流用",
    "AI音声翻訳でブラジル（1.4億）・スペインに自動展開",
    "台湾ユーザー1,225万人（人口の51%）",
    "女性55%・25〜34歳メインがギャルとマッチ",
    "追加制作コストほぼゼロ",
    "フォロワー目標：Year1→5,000、Year3→20万",
]
for i,pt in enumerate(ig_pts):
    txt(s,"▸  "+pt,7.1,2.35+i*0.78,5.6,0.68,size=13,color=DARK)

box(s,6.2,3.6,0.6,0.5)
txt(s,"→",6.25,3.6,0.5,0.5,size=22,bold=True,color=PINK,align=PP_ALIGN.CENTER)

# ═══════ SLIDE 5 インバウンド体験 ════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"インバウンド体験コンテンツ　〜ppgalclub参考〜","渋谷ギャル体験を外国人向け商品化"); num(s,5)

box(s,0.5,1.4,12.3,0.75,bg_color=LPINK,border_color=PINK,bw=1)
txt(s,"参考：@ppgalclub（Instagram 12.1万フォロワー）渋谷でのギャル体験フォトシュート。英語対応・Rakuten Travel Experiencesでも販売済み。",
    0.65,1.5,12.0,0.58,size=13,color=DARK,italic=True)

exp_items = [
    ("体験内容",   "ギャルメイク・衣装着替え＋渋谷フォトシュート（60〜90分）"),
    ("価格設定",   "¥15,000〜30,000 / 人（ppgalclub相場参考）"),
    ("言語対応",   "英語・中国語字幕付き / SNSシェア前提"),
    ("コンテンツ", "体験の様子をTikTok/IGにショートドラマとして投稿 → メディアと連動"),
    ("予約経路",   "Rakuten Travel / Airbnb体験 / DM / ホテルアクティビティ"),
]
for i,(k,v) in enumerate(exp_items):
    top = 2.4+i*0.88
    box(s,0.5,top,3.0,0.72,bg_color=PINK)
    txt(s,k,0.6,top+0.1,2.8,0.5,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    box(s,3.5,top,9.3,0.72,bg_color=LPINK if i%2==0 else WHITE,border_color=PINK,bw=1)
    txt(s,v,3.65,top+0.1,9.0,0.5,size=14,color=DARK)

box(s,0.5,6.85,12.3,0.45,bg_color=ACCENT)
txt(s,"体験→コンテンツ→TikTok配信→新規集客 のサイクルで、撮影コストをメディア制作費として活用",
    0.65,6.88,12.0,0.38,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

# ═══════ SLIDE 6 収益モデル ══════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"収益モデル（5本柱）"); num(s,6)

streams = [
    ("①", "インバウンド体験",  "渋谷ギャル体験",              "¥15K〜30K/人",   "月50→500人"),
    ("②", "広告タイアップ",    "企業×TikTokドラマ",          "200〜400万/案件",  "月1→8件"),
    ("③", "WESELLグッズ販売",  "米国TikTok Shop",             "粗利40〜50%",     "月100→2,000万"),
    ("④", "IPライセンス",       "企業コラボ・公認ロゴ",        "月30〜500万/社",  "2→20社"),
    ("⑤", "ギャル旅タイアップ", "自治体・観光協会",            "300〜700万/地域", "年2→12案件"),
]
header_labels = ["#","収益源","概要","単価","規模感"]
col_widths = [0.5,2.0,3.5,2.5,2.8]
col_starts = [0.4]
for w in col_widths[:-1]: col_starts.append(col_starts[-1]+w+0.05)

for j,(label,cw) in enumerate(zip(header_labels,col_widths)):
    box(s,col_starts[j],1.35,cw,0.5,bg_color=PINK)
    txt(s,label,col_starts[j]+0.05,1.38,cw-0.1,0.4,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

for i,row_data in enumerate(streams):
    bg_c = LPINK if i%2==0 else WHITE
    top = 1.95+i*0.95
    for j,(v,cw) in enumerate(zip(row_data,col_widths)):
        box(s,col_starts[j],top,cw,0.85,bg_color=bg_c,border_color=PINK,bw=1)
        ha = PP_ALIGN.CENTER if j in [0,3,4] else PP_ALIGN.LEFT
        txt(s,v,col_starts[j]+0.08,top+0.08,cw-0.15,0.68,size=13,
            color=ACCENT if j==0 else DARK,bold=(j==0),align=ha)

# ═══════ SLIDE 7 お金の流れ ══════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"お金の流れ・収益分配"); num(s,7)

box(s,0.5,1.35,12.3,0.75,bg_color=LPINK,border_color=PINK,bw=1)
txt(s,"制作費・TikTok運営費はサイバーバズ全額負担　｜　ギャル協会リスクゼロ　｜　ENTIALバズ社内PM（バズの人件費なし）",
    0.65,1.45,12.0,0.58,size=13,bold=True,color=ACCENT,align=PP_ALIGN.CENTER)

nodes = [
    (0.8,  3.2,"企業・自治体\n（広告主）"),
    (4.2,  3.2,"ENTIAL\n（受注・PM）"),
    (7.6,  3.2,"サイバーバズ"),
    (4.9,  5.4,"ギャル協会\nうさたにパイセン"),
    (10.2, 3.2,"WESELL\n（米国EC）"),
]
for l,t,label in nodes:
    box(s,l,t,2.3,1.0,bg_color=PINK,border_color=ACCENT,bw=1.5)
    txt(s,label,l+0.05,t+0.1,2.2,0.8,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

arrows = [
    (2.85,3.6,"広告費→"),
    (6.3,3.6,"売上入金→"),
    (9.65,3.6,"グッズ売上→"),
    (6.8,4.55,"レベニューシェア↓"),
]
for l,t,label in arrows:
    txt(s,label,l,t,2.0,0.45,size=11,color=GRAY,align=PP_ALIGN.CENTER)

shares_info = [
    ("ギャル協会取り分","約32%","出演・監修・IPロイヤリティ"),
    ("サイバーバズ取り分","約50%","運営・制作費・利益"),
    ("ENTIAL取り分","約13%","PM・営業手数料（バズから）"),
]
for i,(name,pct,detail) in enumerate(shares_info):
    top = 1.5+i*0.62 if False else 0
    l = 0.4+i*4.3
    box(s,l,6.05,4.0,1.15,bg_color=LPINK if i%2==0 else WHITE,border_color=PINK,bw=1)
    txt(s,f"{name}　{pct}",l+0.1,6.1,3.8,0.42,size=13,bold=True,color=PINK)
    txt(s,detail,l+0.1,6.5,3.8,0.35,size=11,color=GRAY)

# ═══════ SLIDE 8 収益計画 ════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"3カ年収益計画　Year1：1億　→　Year3：10億"); num(s,8)

years = ["Year1","Year2","Year3"]
data = [
    ("① インバウンド体験", [1800,4800,15000]),
    ("② 広告タイアップ",  [3600,10800,33600]),
    ("③ WESELLグッズ",   [1200,6000,24000]),
    ("④ IPライセンス",    [600,1200,6000]),
    ("⑤ ギャル旅",        [800,2500,6000]),
]
totals = [sum(d[1][i] for d in data) for i in range(3)]

col_w = [3.0,2.7,2.7,2.7]
col_s = [0.4,3.55,6.3,9.05]

for j,(label,cw) in enumerate(zip(["収益源"]+years,col_w)):
    box(s,col_s[j],1.35,cw,0.5,bg_color=PINK)
    txt(s,label,col_s[j]+0.05,1.38,cw-0.1,0.4,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

for i,(name,vals) in enumerate(data):
    top = 1.95+i*0.82
    bg_c = LPINK if i%2==0 else WHITE
    box(s,col_s[0],top,col_w[0],0.72,bg_color=bg_c,border_color=PINK,bw=1)
    txt(s,name,col_s[0]+0.1,top+0.1,col_w[0]-0.15,0.52,size=13,color=DARK)
    for j,v in enumerate(vals):
        box(s,col_s[j+1],top,col_w[j+1],0.72,bg_color=bg_c,border_color=PINK,bw=1)
        txt(s,f"{v:,}万円",col_s[j+1]+0.05,top+0.1,col_w[j+1]-0.1,0.52,
            size=13,color=DARK,align=PP_ALIGN.RIGHT)

top = 1.95+len(data)*0.82
box(s,col_s[0],top,col_w[0],0.72,bg_color=ACCENT)
txt(s,"合　計",col_s[0]+0.1,top+0.1,col_w[0]-0.15,0.52,size=14,bold=True,color=WHITE)
for j,v in enumerate(totals):
    box(s,col_s[j+1],top,col_w[j+1],0.72,bg_color=ACCENT)
    label = f"約{v//10000}億円" if v>=10000 else f"{v:,}万円"
    txt(s,label,col_s[j+1]+0.05,top+0.1,col_w[j+1]-0.1,0.52,
        size=16,bold=True,color=WHITE,align=PP_ALIGN.RIGHT)

# 成長率表示
growth_y = top+0.82
txt(s,"Year1→Year2: +216%成長",col_s[1]+0.1,growth_y,2.5,0.4,size=12,color=ACCENT,bold=True)
txt(s,"Year2→Year3: +234%成長",col_s[2]+0.1,growth_y,2.5,0.4,size=12,color=ACCENT,bold=True)

# ═══════ SLIDE 9 ロードマップ ════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"ロードマップ"); num(s,9)

phases = [
    ("Phase 1\n0〜3ヶ月\n仕込み期",
     ["3社間契約・役割・収益分配を合意",
      "TikTokアカウント開設・初期10本制作",
      "渋谷ギャル体験を商品化（ppgalclub参考）",
      "ENTIAL：営業代理店1〜2社開拓",
      "ホテル・旅行代理店への提案開始"]),
    ("Phase 2\n3〜6ヶ月\n初収益期",
     ["インバウンド体験 月50人達成",
      "タイアップ案件 初回受注",
      "TikTokフォロワー1万人",
      "WESELL稼働・グッズ販売開始",
      "自治体案件1件提案"]),
    ("Phase 3\n6〜36ヶ月\nスケール期",
     ["月次タイアップ案件8件安定稼働",
      "TikTokフォロワー50万人",
      "IG AI翻訳でブラジル・スペイン展開",
      "海外IPライセンス販売開始",
      "年間10億円達成"]),
]
for i,(phase,items) in enumerate(phases):
    l = 0.4+i*4.3
    box(s,l,1.35,4.0,5.8,bg_color=LPINK if i==0 else WHITE,border_color=PINK,bw=2)
    box(s,l,1.35,4.0,1.05,bg_color=PINK)
    txt(s,phase,l+0.1,1.38,3.8,1.0,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    for j,item in enumerate(items):
        txt(s,"□  "+item,l+0.2,2.55+j*0.95,3.6,0.82,size=13,color=DARK)

txt(s,"→",4.25,3.8,0.5,0.5,size=28,bold=True,color=PINK,align=PP_ALIGN.CENTER)
txt(s,"→",8.55,3.8,0.5,0.5,size=28,bold=True,color=PINK,align=PP_ALIGN.CENTER)

# ═══════ SLIDE 10 確認事項 ═══════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s,WHITE); hbar(s,"サイバーバズ打ち合わせ　確認事項"); num(s,10)

items = [
    ("①","TikTokアカウントの名義・権利関係","バズ名義で進める方向で合意確認。WESELLアカウント設計も含む"),
    ("②","ENTIALのバズ社内PM稼働条件","バズの人件費ゼロ前提。ENTIALの報酬はバズからの固定＋インセンティブ"),
    ("③","コンテンツ制作費の上限","Phase1の先行投資額（目安：月200〜300万）"),
    ("④","レベニューシェア比率","ギャル協会：バズ＝4:6をたたき台。ENTIAL分は別途"),
    ("⑤","インバウンド体験の実施体制","ppgalclub参考。渋谷でのギャル体験の運営オペレーション設計"),
    ("⑥","KPI設定","TikTokフォロワー数・月次受注額・体験人数・Year1売上1億円"),
]
for i,(num_,title,detail) in enumerate(items):
    top = 1.5+i*0.95
    box(s,0.4,top,0.7,0.8,bg_color=PINK)
    txt(s,num_,0.42,top+0.12,0.66,0.55,size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    box(s,1.15,top,5.5,0.8,bg_color=LPINK if i%2==0 else WHITE,border_color=PINK,bw=1)
    txt(s,title,1.25,top+0.12,5.3,0.55,size=15,bold=True,color=DARK)
    box(s,6.7,top,6.2,0.8,bg_color=LPINK if i%2==0 else WHITE,border_color=PINK,bw=1)
    txt(s,detail,6.8,top+0.12,6.0,0.55,size=13,color=GRAY,wrap=True)

out = "/home/user/gyal-kyokai--/事業提案書_ギャル×インバウンド_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
