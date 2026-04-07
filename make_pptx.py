from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# カラーパレット（ギャルっぽいピンク系）
PINK       = RGBColor(0xFF, 0x69, 0xB4)   # ホットピンク
LIGHT_PINK = RGBColor(0xFF, 0xD6, 0xE8)   # 薄ピンク
DARK       = RGBColor(0x2D, 0x2D, 0x2D)   # ほぼ黒
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x88, 0x88, 0x88)
ACCENT     = RGBColor(0xFF, 0x14, 0x93)   # ディープピンク

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # 完全blank

# ------------------------------------------------------------------ helpers
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, bg_color=None, border_color=None, border_pt=0):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt2(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as Pt2
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt2(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def header_bar(slide, title, subtitle=None):
    """上部ピンクバー"""
    box(slide, 0, 0, 13.33, 1.2, bg_color=PINK)
    txt(slide, title, 0.4, 0.15, 10, 0.7, size=32, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.82, 10, 0.35, size=14, color=WHITE)

def slide_num(slide, n):
    txt(slide, str(n), 12.8, 7.1, 0.4, 0.3, size=11, color=GRAY, align=PP_ALIGN.RIGHT)

# ================================================================== SLIDE 1 表紙
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_PINK)
box(s, 0, 0, 13.33, 7.5, bg_color=LIGHT_PINK)

# 大きいタイトル背景帯
box(s, 0, 2.2, 13.33, 3.2, bg_color=PINK)

txt(s, "ギャル × インバウンド", 0.5, 2.35, 12.3, 1.1,
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "コンテンツビジネス事業計画書", 0.5, 3.35, 12.3, 0.8,
    size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "ギャル協会  ×  サイバーバズ  ×  ENTIAL", 0.5, 5.7, 12.3, 0.6,
    size=20, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
txt(s, "2026年4月", 0.5, 6.35, 12.3, 0.4,
    size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ================================================================== SLIDE 2 ビジョン
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "ビジョン・事業概要")
slide_num(s, 2)

# 引用ボックス
box(s, 0.5, 1.5, 12.3, 1.3, bg_color=LIGHT_PINK, border_color=PINK, border_pt=2)
txt(s, "「ギャル文化」を日本最強のインバウンドコンテンツIPとして世界に届ける",
    0.8, 1.6, 11.7, 1.0, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(s, "事業概要", 0.5, 3.0, 4, 0.4, size=16, bold=True, color=PINK)

txb = slide.shapes.add_textbox if False else None
tb = s.shapes.add_textbox(Inches(0.5), Inches(3.45), Inches(12.2), Inches(3.5))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "訪日外国人が求める「本物の日本サブカルチャー体験」を、"
r.font.size = Pt(17); r.font.color.rgb = DARK

items = [
    "ギャル協会 / うさたにパイセンのIPをコアに",
    "TikTokショートドラマメディア（サイバーバズ名義）",
    "リアル体験コンテンツ（渋谷ツアー・ギャル旅）",
    "企業・自治体へのタイアップ営業（ENTIAL）",
    "の４軸で収益化するビジネス。"
]
for item in items:
    add_para(tf, ("　▸  " if "の４軸" not in item else "　") + item,
             size=17, color=DARK, space_before=4)

# ================================================================== SLIDE 3 市場背景
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "市場背景　〜なぜ今ギャル×インバウンドか〜")
slide_num(s, 3)

points = [
    ("訪日外国人数が過去最高を更新、インバウンド消費が拡大継続中",),
    ("海外での「ギャル文化」認知がSNSを中心に急上昇（Y2Kブーム・日本サブカル人気）",),
    ("観光庁・自治体がインバウンド向けコンテンツに積極的に予算を投下中",),
    ("ギャル × インバウンドを商品化している競合は現時点でゼロ",),
]

for i, (pt,) in enumerate(points):
    top = 1.5 + i * 1.35
    box(s, 0.5, top, 12.3, 1.1,
        bg_color=LIGHT_PINK if i % 2 == 0 else WHITE,
        border_color=PINK, border_pt=1.5)
    txt(s, f"0{i+1}", 0.65, top+0.18, 0.6, 0.7, size=28, bold=True, color=PINK)
    txt(s, pt, 1.4, top+0.2, 11.0, 0.7, size=17, color=DARK)

# ================================================================== SLIDE 4 3社の役割
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "3社の役割")
slide_num(s, 4)

roles = [
    ("ギャル協会\nうさたにパイセン", "IPオーナー・クリエイター", ["世界観・企画・出演", "ギャル文化の発信源", "コンテンツのコア"]),
    ("サイバーバズ",                "メディアオーナー・実行",   ["TikTokアカウント運営", "広告販売・拡散", "制作費・活動費を負担"]),
    ("ENTIAL",                     "営業部隊・プロジェクト管理", ["企業・自治体への営業", "営業代理店の開拓・管理", "進行・請求管理"]),
]

for i, (name, role, bullets) in enumerate(roles):
    l = 0.4 + i * 4.3
    box(s, l, 1.4, 4.0, 5.6, bg_color=LIGHT_PINK if i==1 else WHITE,
        border_color=PINK, border_pt=2)
    # 名前帯
    box(s, l, 1.4, 4.0, 0.85, bg_color=PINK)
    txt(s, name, l+0.1, 1.42, 3.8, 0.8, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, role, l+0.1, 2.35, 3.8, 0.55, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        txt(s, "▸  " + b, l+0.25, 3.05 + j*0.9, 3.5, 0.75, size=15, color=DARK)

# 矢印テキスト
txt(s, "→", 4.2, 4.1, 0.5, 0.5, size=28, bold=True, color=PINK, align=PP_ALIGN.CENTER)
txt(s, "→", 8.5, 4.1, 0.5, 0.5, size=28, bold=True, color=PINK, align=PP_ALIGN.CENTER)

# ================================================================== SLIDE 5 収益の流れ
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "収益の流れ・お金の設計")
slide_num(s, 5)

# フロー図
nodes = [
    (1.0,  3.2, "企業・自治体\n（広告主）"),
    (4.0,  3.2, "ENTIAL\n（受注窓口）"),
    (7.0,  3.2, "サイバーバズ"),
    (4.7,  5.5, "ギャル協会"),
]
for l, t, label in nodes:
    box(s, l, t, 2.5, 1.1, bg_color=PINK, border_color=ACCENT, border_pt=1)
    txt(s, label, l+0.05, t+0.1, 2.4, 0.9, size=15, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

# 矢印ラベル
arrows = [
    (3.3, 3.55, "広告費・タイアップ費"),
    (6.35, 3.55, "売上入金"),
    (6.8,  4.85, "レベニューシェア"),
    (4.0,  4.5, "営業手数料\n(固定＋インセンティブ)"),
]
for l, t, label in arrows:
    txt(s, "▶  " + label if "入金" in label or "広告" in label else label,
        l, t, 2.8, 0.7, size=12, color=GRAY)

# 補足ノート
box(s, 0.5, 1.4, 12.3, 0.85, bg_color=LIGHT_PINK, border_color=PINK, border_pt=1)
txt(s, "制作費・TikTok運営費はサイバーバズ負担　｜　ギャル協会はリスクゼロ　｜　レベニューシェア比率：バズ：ギャル協会 ＝ 6:4〜7:3（要交渉）",
    0.7, 1.5, 12.0, 0.65, size=13, color=DARK, align=PP_ALIGN.CENTER)

# ================================================================== SLIDE 6 商品パッケージ
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "広告商品パッケージ（インバウンド特化）")
slide_num(s, 6)

pkgs = [
    ("Package 1", "ギャル旅 × 地方観光タイアップ",         "自治体・観光協会",         "300〜700万/地域"),
    ("Package 2", "渋谷ギャル文化体験ツアー × メディア化",  "ホテル・旅行代理店",        "50〜150万"),
    ("Package 3", "インバウンド向けショートドラマ広告",      "免税店・コスメ・決済系",    "150〜400万"),
    ("Package 4", "「ギャル協会公認」コンテンツライセンス",  "海外メディア・観光アプリ",  "月20〜80万"),
]

for i, (pkg, name, target, price) in enumerate(pkgs):
    top = 1.45 + i * 1.42
    box(s, 0.4, top, 1.5, 1.2, bg_color=PINK)
    txt(s, pkg, 0.42, top+0.25, 1.46, 0.7, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.9, top, 6.5, 1.2, bg_color=LIGHT_PINK if i%2==0 else WHITE,
        border_color=PINK, border_pt=1)
    txt(s, name, 2.05, top+0.28, 6.2, 0.65, size=16, bold=True, color=DARK)
    box(s, 8.4, top, 2.7, 1.2,
        bg_color=LIGHT_PINK if i%2==0 else WHITE, border_color=PINK, border_pt=1)
    txt(s, target, 8.5, top+0.28, 2.5, 0.65, size=14, color=DARK)
    box(s, 11.1, top, 2.1, 1.2, bg_color=PINK)
    txt(s, price, 11.15, top+0.28, 2.0, 0.65, size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

# ヘッダー行
for l, w, label in [(1.9,6.5,"商品名"), (8.4,2.7,"ターゲット"), (11.1,2.1,"参考単価")]:
    txt(s, label, l+0.1, 1.05, w-0.1, 0.35, size=12, bold=True, color=GRAY)

# ================================================================== SLIDE 7 営業戦略
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "営業戦略（ENTIALの動き方）")
slide_num(s, 7)

txt(s, "ターゲット優先順位", 0.5, 1.4, 5, 0.4, size=16, bold=True, color=PINK)

targets = [
    ("★★★", "ホテル・旅行代理店",   "意思決定が速く予算あり。インバウンド需要でピーク"),
    ("★★★", "免税店・百貨店",       "インバウンド消費に直結。単価が高く継続案件になりやすい"),
    ("★★",  "地方自治体・観光協会", "予算大だが意思決定遅い。補助金活用でハードル下げられる"),
    ("★",   "OTA・体験予約サービス","拡散力大。TikTokが育ってから本格攻略"),
]
for i, (star, name, reason) in enumerate(targets):
    top = 1.9 + i * 1.1
    box(s, 0.5, top, 1.0, 0.85, bg_color=LIGHT_PINK, border_color=PINK, border_pt=1)
    txt(s, star, 0.5, top+0.12, 1.0, 0.6, size=14, color=PINK, align=PP_ALIGN.CENTER)
    txt(s, name, 1.65, top+0.12, 3.3, 0.6, size=15, bold=True, color=DARK)
    txt(s, reason, 5.1, top+0.12, 8.0, 0.6, size=14, color=GRAY)

txt(s, "ENTIALの動き方", 0.5, 6.25, 5, 0.4, size=16, bold=True, color=PINK)
steps = "①営業代理店1〜2社を開拓（インバウンド・観光系に強い代理店）　→　②代理店経由でホテル・旅行代理店・百貨店へ提案　→　③自治体案件は補助金スキームとセットで直接提案"
txt(s, steps, 0.5, 6.65, 12.5, 0.65, size=14, color=DARK)

# ================================================================== SLIDE 8 ロードマップ
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "ロードマップ")
slide_num(s, 8)

phases = [
    ("Phase 1\n0〜3ヶ月\n仕込み期",
     ["3社間の契約・役割・収益分配を合意",
      "TikTokアカウント開設・初期コンテンツ10本制作",
      "ENTIALが営業代理店1〜2社を開拓",
      "ホテル・旅行代理店への初回提案開始"]),
    ("Phase 2\n3〜6ヶ月\n初収益期",
     ["タイアップ案件 初回受注",
      "TikTokフォロワー 1万人達成",
      "自治体案件 1件提案・交渉開始"]),
    ("Phase 3\n6〜12ヶ月\nスケール期",
     ["月次タイアップ案件を安定受注",
      "TikTok英語字幕で海外フォロワー獲得",
      "ギャルIPの海外ライセンス販売開始",
      "神7組成・複数案件を並走できる体制"]),
]

for i, (phase, items) in enumerate(phases):
    l = 0.4 + i * 4.3
    box(s, l, 1.35, 4.0, 5.7,
        bg_color=LIGHT_PINK if i==0 else WHITE, border_color=PINK, border_pt=2)
    box(s, l, 1.35, 4.0, 1.1, bg_color=PINK)
    txt(s, phase, l+0.1, 1.38, 3.8, 1.05, size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txt(s, "□  " + item, l+0.2, 2.6 + j*1.0, 3.6, 0.85, size=14, color=DARK)

txt(s, "→", 4.25, 4.0, 0.5, 0.5, size=30, bold=True, color=PINK, align=PP_ALIGN.CENTER)
txt(s, "→", 8.55, 4.0, 0.5, 0.5, size=30, bold=True, color=PINK, align=PP_ALIGN.CENTER)

# ================================================================== SLIDE 9 確認事項
s = prs.slides.add_slide(blank_layout)
bg(s, WHITE)
header_bar(s, "サイバーバズ打ち合わせ　確認事項")
slide_num(s, 9)

items = [
    ("①", "TikTokアカウントの名義・権利関係",    "バズ名義で進める方向で合意確認"),
    ("②", "コンテンツ制作費の上限設定",          "Phase 1の先行投資額を決める"),
    ("③", "レベニューシェア比率",                "ギャル協会との分配。6:4をたたき台に交渉"),
    ("④", "ENTIALへの営業手数料スキーム",        "固定費＋売上連動インセンティブの設計"),
    ("⑤", "KPI設定",                            "TikTokフォロワー数・月次受注額・1年の売上目標"),
]
for i, (num, title, detail) in enumerate(items):
    top = 1.5 + i * 1.1
    box(s, 0.4, top, 0.65, 0.85, bg_color=PINK)
    txt(s, num, 0.4, top+0.12, 0.65, 0.6, size=18, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.2, top+0.05, 5.5, 0.45, size=16, bold=True, color=DARK)
    txt(s, detail, 1.2, top+0.48, 11.0, 0.4, size=14, color=GRAY)

# ================================================================== save
out = "/home/user/gyal-kyokai--/事業計画書_ギャル×インバウンド.pptx"
prs.save(out)
print(f"Saved: {out}")
